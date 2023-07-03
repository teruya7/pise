from pptx.util import Cm, Pt
from pdf2image import convert_from_path
from pathlib import Path
import os
import json

def add_slide(presentation_object, layout_number):
    layout_object = presentation_object.slide_layouts[layout_number]
    slide = presentation_object.slides.add_slide(layout_object)
    return slide

def add_image(slide, picture_path, left=Cm(5), top=Cm(5), width=Cm(15), height=None):
    shapes = slide.shapes
    shapes.add_picture(picture_path, left, top, width, height)

def add_title(slide,title):
    title_box = slide.placeholders[0] 
    title_box.text = title 

def add_data(slide, text_data, left=Cm(10), top=Cm(3), width=Cm(15), height=Cm(15)):
    text_box = slide.shapes.add_textbox(left, top, width, height) 
    frame = text_box.text_frame
    paragraph = frame.add_paragraph()
    paragraph.text = text_data
    paragraph.font.size = Pt(15)

def read_data(data_file):
    with open(data_file) as f:
        data = f.read()
    return data

def make_slide_with_data(prs, data_name, path):
    slide = add_slide(prs, 1)
    add_title(slide,data_name)
    os.chdir(path)
    try:
        data = read_data(data_name)
        add_data(slide, data)
    except FileNotFoundError:
        print(f'No {data_name}')

def pdf_to_png(pdf_file, img_path, fmt='png', dpi=200):

    #pdf_file、img_pathをPathにする
    pdf_path = Path(pdf_file)
    image_dir = Path(img_path)

    # PDFをImage に変換(pdf2imageの関数)
    pages = convert_from_path(pdf_path, dpi)

    # 画像ファイルを１ページずつ保存
    for i, page in enumerate(pages):
        file_name = "{}_{:02d}.{}".format(pdf_path.stem,i+1,fmt)
        image_path = image_dir / file_name
        page.save(image_path, fmt)
        before_name = image_path
        after_name = pdf_path.stem + "." + fmt
        os.rename(before_name, after_name)

def make_slide_with_image(prs, image_name, path):
    slide = add_slide(prs, 1)
    add_title(slide, image_name)
    os.chdir(path)
    if os.path.isfile(f"{image_name}.pdf"):
        pdf_to_png(f"./{image_name}.pdf", "./")
        add_image(slide, f"./{image_name}.png")
    elif os.path.isfile(f"{image_name}.png"):
        add_image(slide, f"./{image_name}.png")
    else:
        print(f'No {image_name} file exists')
    return slide
        
#save_pathにname.pptxとして保存
def save_pptx(prs, name, save_path):
    os.chdir(save_path)
    prs.save(f"{name}.pptx")

#化学ポテンシャルの極限のラベルを取得
def get_label_from_chempotdiag(path_chem_pot_diag):
    label = []
    with open(path_chem_pot_diag) as f:
        chem_pot = json.load(f)
    for label_key in chem_pot["target_vertices_dict"]:
        label.append(label_key)
    return label

def check_analysis_info(path):
    with open(path) as f:
        analysis_info = json.load(f)
        for boolean in analysis_info.values():
            if boolean:
                flag = True
            else:
                flag = False
                break
    return flag



